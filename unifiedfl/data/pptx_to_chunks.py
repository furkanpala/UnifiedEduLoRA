"""
Convert a PowerPoint (.pptx) file into plain-text chunks suitable for generate_qa.py.

BEFORE YOU RUN THIS — things you will likely need to adjust:
  - SKIP_LAYOUTS    : slide layout names to skip (title slides, section dividers, etc.)
                      Print layout names first with --list-layouts to see what yours are called.
  - SKIP_IF_FEWER   : minimum words a slide must contain to be included. Slides with only
                      a title and a single bullet are often too thin for good QA — raise this
                      if you find many low-quality chunks.
  - TARGET_WORDS    : target chunk size. PowerPoint slides tend to be sparse, so the default
                      merges several slides before cutting. Increase for text-heavy decks.
  - INCLUDE_TITLES  : set to False if slide titles are just labels (e.g. "Slide 3") and
                      add no semantic value.
  - The note extractor is OFF by default (INCLUDE_NOTES = False). Turn it on if your slides
                      have speaker notes that contain the real explanation.

Requirements:
  pip install python-pptx

Usage:
  python pptx_to_chunks.py lecture.pptx
  python pptx_to_chunks.py lecture.pptx --output my_chunks.json
  python pptx_to_chunks.py lecture.pptx --list-layouts        # inspect layout names first
  python pptx_to_chunks.py lecture.pptx --target-words 300 --skip-if-fewer 30
"""

from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import List, Set


# ── Adjust these to match your presentation ──────────────────────────────────

# Layout names (case-insensitive substrings) whose slides are silently dropped.
# Common values: "Title Slide", "Section Header", "Blank", "Picture with Caption"
SKIP_LAYOUTS: List[str] = ["title slide", "section header", "blank"]

SKIP_IF_FEWER  = 20    # slides with fewer body words than this are dropped
TARGET_WORDS   = 250   # target words per output chunk
MIN_WORDS      = 150   # chunk must have at least this many words
MAX_WORDS      = 400   # hard upper limit per chunk

INCLUDE_TITLES = True  # include slide titles in the extracted text
INCLUDE_NOTES  = False # include speaker notes (set True if notes contain explanations)


# ── Text extraction ───────────────────────────────────────────────────────────

def _slide_text(slide, include_titles: bool, include_notes: bool) -> str:
    from pptx.util import Pt
    from pptx.enum.text import PP_ALIGN

    lines: List[str] = []

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        # Detect title placeholder
        is_title = (
            hasattr(shape, "placeholder_format")
            and shape.placeholder_format is not None
            and shape.placeholder_format.idx in (0, 1)
        )

        if is_title and not include_titles:
            continue

        for para in shape.text_frame.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            # Skip lines that are purely numeric (slide numbers, etc.)
            if re.fullmatch(r"\d+", text):
                continue
            lines.append(text)

    if include_notes and slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        for para in notes_tf.paragraphs:
            text = para.text.strip()
            if text:
                lines.append(text)

    return " ".join(lines)


def _clean(text: str) -> str:
    # Collapse runs of whitespace
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n+", " ", text)
    return text.strip()


# ── Chunking ──────────────────────────────────────────────────────────────────

def _split_into_chunks(text: str, target: int, min_w: int, max_w: int) -> List[str]:
    words = text.split()
    chunks: List[str] = []
    i = 0
    while i < len(words):
        end = min(i + target, len(words))
        candidate = " ".join(words[i:end])
        if end < len(words):
            for boundary in range(end, max(i + min_w, end - 30), -1):
                snippet = " ".join(words[i:boundary])
                if re.search(r"[.!?]\s*$", snippet):
                    candidate = snippet
                    end = boundary
                    break
        chunk = candidate.strip()
        chunk_words = chunk.split()
        while len(chunk_words) > max_w:
            chunks.append(" ".join(chunk_words[:max_w]))
            chunk_words = chunk_words[max_w:]
        if len(chunk_words) >= min_w:
            chunks.append(" ".join(chunk_words))
        elif chunks:
            chunks[-1] += " " + " ".join(chunk_words)
        i = end
    return chunks


# ── Main extraction ───────────────────────────────────────────────────────────

def pptx_to_chunks(
    pptx_path: Path,
    skip_layouts: List[str] = SKIP_LAYOUTS,
    skip_if_fewer: int = SKIP_IF_FEWER,
    target_words: int = TARGET_WORDS,
    min_words: int = MIN_WORDS,
    max_words: int = MAX_WORDS,
    include_titles: bool = INCLUDE_TITLES,
    include_notes: bool = INCLUDE_NOTES,
) -> List[str]:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")

    prs = Presentation(str(pptx_path))
    skip_lower: Set[str] = {s.lower() for s in skip_layouts}

    slide_texts: List[str] = []
    for slide in prs.slides:
        layout_name = slide.slide_layout.name.lower()
        if any(pattern in layout_name for pattern in skip_lower):
            continue

        text = _clean(_slide_text(slide, include_titles, include_notes))
        word_count = len(text.split())

        if word_count < skip_if_fewer:
            continue

        slide_texts.append(text)

    full_text = " ".join(slide_texts)
    return _split_into_chunks(full_text, target_words, min_words, max_words)


def list_layouts(pptx_path: Path) -> None:
    try:
        from pptx import Presentation
    except ImportError:
        raise ImportError("pip install python-pptx")

    prs = Presentation(str(pptx_path))
    layouts_seen: dict = {}
    for slide in prs.slides:
        name = slide.slide_layout.name
        layouts_seen[name] = layouts_seen.get(name, 0) + 1

    print("Slide layout names and counts:")
    for name, count in sorted(layouts_seen.items(), key=lambda x: -x[1]):
        print(f"  {count:>3}x  {name!r}")
    print("\nAdd unwanted layout names to SKIP_LAYOUTS in this script.")


# ── CLI ───────────────────────────────────────────────────────────────────────

def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Convert a PPTX file to plain-text chunks for generate_qa.py"
    )
    p.add_argument("pptx", help="Path to the .pptx file")
    p.add_argument("--output", default=None,
                   help="Output JSON file (default: <pptx_stem>_chunks.json)")
    p.add_argument("--list-layouts", action="store_true",
                   help="Print slide layout names and exit (useful to configure SKIP_LAYOUTS)")
    p.add_argument("--skip-if-fewer", type=int, default=SKIP_IF_FEWER,
                   help=f"Drop slides with fewer body words than this (default: {SKIP_IF_FEWER})")
    p.add_argument("--target-words", type=int, default=TARGET_WORDS,
                   help=f"Target words per chunk (default: {TARGET_WORDS})")
    p.add_argument("--include-notes", action="store_true", default=INCLUDE_NOTES,
                   help="Also extract speaker notes from each slide")
    return p.parse_args()


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx)

    if not pptx_path.exists():
        print(f"ERROR: file not found: {pptx_path}", file=sys.stderr)
        sys.exit(1)

    if args.list_layouts:
        list_layouts(pptx_path)
        return

    print(f"Reading {pptx_path} ...")
    chunks = pptx_to_chunks(
        pptx_path,
        skip_if_fewer=args.skip_if_fewer,
        target_words=args.target_words,
        include_notes=args.include_notes,
    )

    if not chunks:
        print("WARNING: no chunks produced. Try lowering --skip-if-fewer or checking SKIP_LAYOUTS.")
        sys.exit(1)

    out_path = Path(args.output) if args.output else pptx_path.with_name(pptx_path.stem + "_chunks.json")
    out_path.write_text(json.dumps(chunks, indent=2, ensure_ascii=False), encoding="utf-8")

    print(f"Extracted {len(chunks)} chunks → {out_path}")
    print(f"Word counts: min={min(len(c.split()) for c in chunks)}  "
          f"max={max(len(c.split()) for c in chunks)}  "
          f"avg={sum(len(c.split()) for c in chunks)//len(chunks)}")
    print("\nNext step: feed these chunks into generate_qa.py")
    print("  my_chunks = json.load(open(\"" + str(out_path) + "\"))")


if __name__ == "__main__":
    main()
